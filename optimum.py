from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def add_content_slide(prs, title, bullet_points, illustration_note=""):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear existing
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.level = 0
        p.space_before = Pt(6)
        
    # Add illustration note at the bottom
    if illustration_note:
        p = tf.add_paragraph()
        p.text = illustration_note
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.space_before = Pt(20)

def create_presentation():
    prs = Presentation()

    # SLIDE 1: Title
    add_title_slide(prs, 
        "Transformation de notre modèle de production digitale", 
        "Vers une approche 'Plug & Play' pour accélérer la mise sur le marché.\nPrésenté par : [Votre Nom/Équipe]")

    # SLIDE 2: Le Constat
    add_content_slide(prs, 
        "Notre Défi Actuel : La duplication d'efforts",
        [
            "Le problème : Nous redéveloppons les mêmes fonctionnalités pour chaque partenaire.",
            "Conséquence 1 : Délais de livraison longs pour des besoins standards.",
            "Conséquence 2 : Coûts de développement gonflés.",
            "Conséquence 3 : Risques de bugs récurrents et maintenance difficile."
        ],
        "💡 Illustration : Graphique montrant une courbe de coûts croissante à chaque nouveau client.")

    # SLIDE 3: La Solution
    add_content_slide(prs,
        "La Solution : L'Approche Modulaire",
        [
            "Concept : Une architecture en 'Briques Lego'.",
            "Objectif : Assembler des modules existants plutôt que de recoder.",
            "Bénéfice : Rapidité, Économie et Fiabilité."
        ],
        "💡 Illustration : Image de blocs Lego assemblés (Bloc Auth, Bloc RH, Bloc Courriers).")

    # SLIDE 4: Catalogue des Modules
    add_content_slide(prs,
        "Notre Catalogue de Modules Réutilisables",
        [
            "🕒 Pointage & Gestion de Présence",
            "🏖️ Gestion des Congés & Permissions",
            "📅 Agenda & Planification",
            "📬 Gestion des Courriers",
            "🚶 Gestion des Visites",
            "📝 Notes de Services",
            "👥 Ressources Humaines (RH)",
            "🔐 Authentification & Sécurité"
        ],
        "💡 Illustration : Une interface montrant un catalogue avec des boutons 'Activer/Désactiver'.")

    # SLIDE 5: Architecture Simplifiée
    add_content_slide(prs,
        "Architecture Technique Simplifiée",
        [
            "Couche Présentation : Interface adaptée au logo/couleurs du client.",
            "Couche Modules : Le moteur. Les fonctionnalités s'activent à la demande.",
            "Couche Données : Isolation stricte. Chaque client possède son espace sécurisé.",
            "Cloud Dédicé : Possibilité d'utiliser le compte AWS/Cloudinary du client."
        ],
        "💡 Illustration : Schéma en 3 niveaux (Interface -> Moteur -> Coffre-fort).")

    # SLIDE 6: Sécurité (Point Clé CEO)
    add_content_slide(prs,
        "Sécurité & Souveraineté des Données",
        [
            "Isolation des données (Multi-tenancy) : Aucun mélange entre clients.",
            "Souveraineté : Option 'Bring Your Own Cloud' (BYOC).",
            "Le client garde la propriété juridique et financière de ses fichiers.",
            "Arguments de vente forts pour les grands comptes (Conformité RGPD)."
        ],
        "💡 Illustration : Image de silos étanches ou de cadenas individuels.")

    # SLIDE 7: Avantages vs Inconvénients
    slide_layout = prs.slide_layouts[5] # Blank for table
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_shape.text = "Analyse : Avantages et Inconvénients"
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Table
    rows, cols = 4, 2
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    table.cell(0, 0).text = "✅ AVANTAGES"
    table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
    table.cell(0, 1).text = "⚠️ RISQUES / INCONVÉNIENTS"
    table.cell(0, 1).text_frame.paragraphs[0].font.bold = True
    
    # Content
    data = [
        ("Vitesse : Déploiement en jours, pas en mois.", "Investissement initial élevé (R&D)."),
        ("Coûts réduits : Marges augmentées.", "Rigidité : Personnalisation spécifique plus complexe."),
        ("Qualité : Code éprouvé et testé par tous.", "Maintenance globale : Une mise à jour impacte tous les clients."),
        ("Sécurité : Isolation des données rassurante.", "Compétences techniques requises pour l'équipe.")
    ]
    
    for i, (av, inco) in enumerate(data):
        table.cell(i+1, 0).text = av
        table.cell(i+1, 1).text = inco

    # SLIDE 8: Conclusion
    add_content_slide(prs,
        "Conclusion & Impact Business",
        [
            "Passage du statut 'Prestataire' à 'Éditeur de Solution'.",
            "ROI : Rentabilité accrue sur le long terme.",
            "Avantage Concurrentiel : Réactivité et Sécurité.",
            "Prochaine étape : Lancer le premier pilote avec les modules Auth et RH."
        ],
        "💡 Illustration : Graphique montrant la rentabilité future vs modèle actuel.")

    # Save
    prs.save('Presentation_Architecture_Digitale.pptx')
    print("✅ Fichier PowerPoint 'Presentation_Architecture_Digitale.pptx' généré avec succès !")

if __name__ == "__main__":
    create_presentation()